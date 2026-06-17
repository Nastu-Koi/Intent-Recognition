"""
General Chat 工具集 — 图像识别与文档总结。

每个工具内部完成「上传文件到 Dify → 调用 Dify API」的完整流程，
调用方无需关心 Dify file_id 的获取。

工具列表:
  1. image_recognition  — 图片 OCR / 场景分析 / 发票识别
  2. document_summary   — 文档总结 / 要点提炼
"""

import csv
import os
import re
import time
from io import BytesIO, StringIO
from pathlib import Path
from typing import List

import requests
from langchain_core.tools import tool

EMU_PER_INCH = 914400

from engine.logging_config import get_logger

logger = get_logger(__name__)


# ──────────────────────────────────────────────
# Dify 文件上传（内部共享函数）
# ──────────────────────────────────────────────

def _get_dify_base_url() -> str:
    return (
        os.getenv("DIFY_API_BASE_URL")
        or os.getenv("DIFY_BASE_URL")
        or "https://api.dify.ai/v1"
    ).rstrip("/")


def _get_dify_api_key(agent_id: str = "") -> str:
    """
    按优先级获取 Dify API Key:
      1. DIFY_<AGENT_ID>_API_KEY  (例: DIFY_VISION_API_KEY)
      2. DIFY_API_KEY             (全局)
    """
    if agent_id:
        normalized = agent_id.upper().replace("-", "_")
        key = os.getenv(f"DIFY_{normalized}_API_KEY", "")
        if key:
            return key
    return os.getenv("DIFY_API_KEY", "")


def _guess_mime(path: Path) -> str:
    """根据扩展名猜测 MIME 类型。"""
    ext = path.suffix.lower()
    mime_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".bmp": "image/bmp",
        ".pdf": "application/pdf",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".csv": "text/csv",
    }
    return mime_map.get(ext, "application/octet-stream")


# 统一的 Dify 用户标识 (上传与调用必须一致)
DIFY_USER = os.getenv("DIFY_USER", "intent-recognition")


def _upload_file_to_dify(file_path: str, api_key: str = "") -> str:
    """
    上传本地文件到 Dify，返回 file_id。

    重要: api_key 必须是目标 Dify App 的 API Key，否则上传的 file_id
    在后续调用该 App 的 chat/workflow 时会找不到文件。

    Args:
        file_path: 本地文件绝对路径
        api_key: 目标 Dify App 的 API Key (必须与后续调用 App 的 key 一致)

    Returns:
        Dify file_id 字符串

    Raises:
        FileNotFoundError: 文件不存在
        ValueError: API Key 未配置
        requests.HTTPError: 上传失败
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    key = api_key or _get_dify_api_key()
    if not key:
        raise ValueError("DIFY_API_KEY 未配置，无法上传文件。")

    url = f"{_get_dify_base_url()}/files/upload"
    headers = {"Authorization": f"Bearer {key}"}

    with open(path, "rb") as f:
        files = {"file": (path.name, f, _guess_mime(path))}
        data = {"user": DIFY_USER}
        resp = requests.post(url, headers=headers, files=files, data=data, timeout=120)
        resp.raise_for_status()
        result = resp.json()

    file_id = result.get("id") or result.get("file_id", "")
    logger.info(f"[Tools] 文件上传成功: {path.name} -> file_id={file_id}")
    return file_id


# ──────────────────────────────────────────────
# Dify Vision 调用（内部函数）
# ──────────────────────────────────────────────

def _call_dify_vision(
    query: str,
    file_ids: List[str],
    app_type: str = "chat",
    conversation_id: str = "",
) -> str:
    """
    调用 Dify Vision App（支持 chat / workflow 模式）。

    Args:
        query: 识别指令
        file_ids: Dify file_id 列表
        app_type: chat 或 workflow
        conversation_id: Dify 会话 ID（用于多轮对话）

    Returns:
        识别结果文本
    """
    api_key = _get_dify_api_key("VISION")
    if not api_key:
        raise ValueError("DIFY_VISION_API_KEY 或 DIFY_API_KEY 未配置。")

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    files_payload = [
        {
            "type": "image",
            "transfer_method": "local_file",
            "upload_file_id": fid,
        }
        for fid in file_ids
    ]

    if app_type == "workflow":
        url = f"{_get_dify_base_url()}/workflows/run"
        payload = {
            "inputs": {"query": query},
            "response_mode": "blocking",
            "user": DIFY_USER,
            "files": files_payload,
        }
        if conversation_id:
            payload["conversation_id"] = conversation_id
        resp = requests.post(url, json=payload, headers=headers, timeout=120)
        resp.raise_for_status()
        data = resp.json()
        outputs = (data.get("data") or {}).get("outputs") or data.get("outputs") or {}
        if isinstance(outputs, dict):
            for key in ("answer", "result", "text", "output"):
                if key in outputs:
                    return str(outputs[key])
        return str(outputs or data)
    else:  # chat
        url = f"{_get_dify_base_url()}/chat-messages"
        payload = {
            "inputs": {},
            "query": query,
            "response_mode": "blocking",
            "conversation_id": conversation_id,
            "user": DIFY_USER,
            "files": files_payload,
        }
        resp = requests.post(url, json=payload, headers=headers, timeout=120)
        resp.raise_for_status()
        data = resp.json()
        return data.get("answer") or data.get("text") or str(data)


# ──────────────────────────────────────────────
# Dify Doc Summary 调用（内部函数）
# ──────────────────────────────────────────────

def _call_dify_doc_summary(
    query: str,
    file_ids: List[str],
    app_type: str = "chat",
    conversation_id: str = "",
) -> str:
    """
    调用 Dify 文档总结 App。

    Args:
        query: 总结指令
        file_ids: Dify file_id 列表
        app_type: chat 或 workflow
        conversation_id: Dify 会话 ID（用于多轮对话）

    Returns:
        总结结果文本
    """
    api_key = _get_dify_api_key("DOC_SUMMARY")
    if not api_key:
        raise ValueError("DIFY_DOC_SUMMARY_API_KEY 或 DIFY_API_KEY 未配置。")

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    # 构建 files 参数
    files_payload = [
        {
            "type": "document",
            "transfer_method": "local_file",
            "upload_file_id": fid,
        }
        for fid in file_ids
    ]

    if app_type == "workflow":
        input_variable = os.getenv("DIFY_DOC_SUMMARY_INPUT_VARIABLE", "query")
        url = f"{_get_dify_base_url()}/workflows/run"
        payload = {
            "inputs": {input_variable: query},
            "response_mode": "blocking",
            "user": DIFY_USER,
            "files": files_payload,
        }
        if conversation_id:
            payload["conversation_id"] = conversation_id
        resp = requests.post(url, json=payload, headers=headers, timeout=120)
        resp.raise_for_status()
        data = resp.json()
        outputs = (data.get("data") or {}).get("outputs") or data.get("outputs") or {}
        if isinstance(outputs, dict):
            for key in ("answer", "result", "text", "output", "summary", "summarized"):
                if key in outputs:
                    return str(outputs[key])
        return str(outputs or data)
    else:  # chat
        url = f"{_get_dify_base_url()}/chat-messages"
        payload = {
            "inputs": {},
            "query": query,
            "response_mode": "blocking",
            "conversation_id": conversation_id,
            "user": DIFY_USER,
            "files": files_payload,
        }
        resp = requests.post(url, json=payload, headers=headers, timeout=120)
        resp.raise_for_status()
        data = resp.json()
        return data.get("answer") or data.get("text") or str(data)


# ══════════════════════════════════════════════
# LangChain Tool 定义
# ══════════════════════════════════════════════

# 支持的图片扩展名
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
# 支持的文档扩展名
_DOCUMENT_EXTENSIONS = {".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".txt", ".md", ".csv"}


def _safe_output_filename(filename: str, default_stem: str, suffix: str) -> str:
    name = Path(filename or "").name.strip()
    if not name:
        name = f"{default_stem}_{int(time.time())}{suffix}"
    if Path(name).suffix.lower() != suffix:
        name = f"{Path(name).stem or default_stem}{suffix}"
    stem = re.sub(r"[^\w\-.一-鿿]+", "_", Path(name).stem, flags=re.UNICODE).strip("._")
    return f"{stem or default_stem}{suffix}"


def _get_upload_dir() -> Path:
    upload_dir = Path(__file__).resolve().parents[2] / "uploads"
    upload_dir.mkdir(exist_ok=True)
    return upload_dir


def _parse_slide_sections(content: str) -> List[dict]:
    slides = []
    current = None
    for raw_line in (content or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue

        if line.startswith("#"):
            title = line.lstrip("#").strip()
            if title:
                current = {"title": title, "bullets": []}
                slides.append(current)
            continue

        if current is None:
            current = {"title": line, "bullets": []}
            slides.append(current)
        elif line.startswith(("- ", "* ")):
            current["bullets"].append(line[2:].strip())
        elif re.match(r"^\d+[\.)]\s+", line):
            current["bullets"].append(re.sub(r"^\d+[\.)]\s+", "", line).strip())
        else:
            current["bullets"].append(line)

    return slides


def _parse_table_content(content: str) -> List[List[str]]:
    text = (content or "").strip()
    if not text:
        return []

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return []

    if any("|" in line for line in lines):
        rows = []
        for line in lines:
            if re.fullmatch(r"[|:\-\s]+", line):
                continue
            cells = [cell.strip() for cell in line.strip("|").split("|")]
            if cells:
                rows.append(cells)
        return rows

    try:
        delimiter = "\t" if any("\t" in line for line in lines) else ","
        return [[cell.strip() for cell in row] for row in csv.reader(StringIO(text), delimiter=delimiter) if row]
    except csv.Error:
        return [[line] for line in lines]


@tool
def image_recognition(file_path: str, instruction: str = "请描述这张图片的内容") -> str:
    """识别和分析图片内容。支持 OCR 文字识别、发票识别、场景分析、图片内容描述等视觉任务。

    Args:
        file_path: 图片文件的本地路径（必须是服务器上的绝对路径）
        instruction: 识别指令，例如「请识别这张发票上的金额」、「描述图片内容」、「提取图中的文字」
    """
    logger.info(f"[Tool:image_recognition] file={file_path} | instruction={instruction[:100]}")

    # 文件类型校验：仅接受图片文件
    ext = Path(file_path).suffix.lower()
    if ext not in _IMAGE_EXTENSIONS:
        return (
            f"错误: 文件 {Path(file_path).name} 不是支持的图片格式 "
            f"(支持: {', '.join(_IMAGE_EXTENSIONS)})。"
            f"如需处理文档，请使用 document_summary 工具。"
        )

    try:
        # 获取目标 App 的 API Key（上传和调用必须使用同一个 key）
        vision_api_key = _get_dify_api_key("VISION")

        # Step 1: 使用 Vision App 的 key 上传图片到 Dify
        file_id = _upload_file_to_dify(file_path, api_key=vision_api_key)

        # Step 2: 调用 Dify Vision API
        app_type = os.getenv("DIFY_VISION_APP_TYPE", "chat")
        result = _call_dify_vision(
            query=instruction,
            file_ids=[file_id],
            app_type=app_type,
        )

        logger.info(f"[Tool:image_recognition] 识别完成，结果长度={len(result)}")
        return result

    except FileNotFoundError as e:
        return f"错误: {e}"
    except Exception as e:
        logger.error(f"[Tool:image_recognition] 执行失败: {e}")
        return f"图片识别失败: {e}"


@tool
def document_summary(file_path: str, instruction: str = "请总结这份文档的核心内容") -> str:
    """对文档进行智能总结和要点提炼。支持 PDF、Word、Excel、TXT、Markdown 等格式。

    Args:
        file_path: 文档文件的本地路径（必须是服务器上的绝对路径）
        instruction: 总结指令，例如「总结核心要点」、「提炼关键条款」、「分析报告结论」
    """
    logger.info(f"[Tool:document_summary] file={file_path} | instruction={instruction[:100]}")

    # 文件类型校验：仅接受文档文件
    ext = Path(file_path).suffix.lower()
    if ext not in _DOCUMENT_EXTENSIONS:
        return (
            f"错误: 文件 {Path(file_path).name} 不是支持的文档格式 "
            f"(支持: {', '.join(_DOCUMENT_EXTENSIONS)})。"
            f"如需处理图片，请使用 image_recognition 工具。"
        )

    try:
        # 获取目标 App 的 API Key（上传和调用必须使用同一个 key）
        doc_summary_api_key = _get_dify_api_key("DOC_SUMMARY")

        # Step 1: 使用 Doc Summary App 的 key 上传文档到 Dify
        file_id = _upload_file_to_dify(file_path, api_key=doc_summary_api_key)

        # Step 2: 调用 Dify 文档总结 API
        app_type = os.getenv("DIFY_DOC_SUMMARY_APP_TYPE", "chat")
        result = _call_dify_doc_summary(
            query=instruction,
            file_ids=[file_id],
            app_type=app_type,
        )

        logger.info(f"[Tool:document_summary] 总结完成，结果长度={len(result)}")
        return result

    except FileNotFoundError as e:
        return f"错误: {e}"
    except Exception as e:
        logger.error(f"[Tool:document_summary] 执行失败: {e}")
        return f"文档总结失败: {e}"


@tool
def pdf_add_watermark(
    file_path: str,
    watermark_text: str = "WATERMARK",
    opacity: float = 0.18,
    font_size: int = 48,
    rotation: int = 45,
) -> str:
    """给 PDF 的每一页添加文本水印，并生成新的 PDF 文件。

    Args:
        file_path: PDF 文件的本地路径（必须是服务器上的绝对路径）
        watermark_text: 要添加的水印文字，例如「CONFIDENTIAL」或「仅供内部使用」
        opacity: 水印透明度，0 到 1 之间，默认 0.18
        font_size: 水印字号，默认 48
        rotation: 水印旋转角度，默认 45
    """
    logger.info(
        "[Tool:pdf_add_watermark] file=%s | text=%s | opacity=%s | font_size=%s | rotation=%s",
        file_path,
        watermark_text,
        opacity,
        font_size,
        rotation,
    )

    path = Path(file_path)
    if not path.exists():
        return f"错误: 文件不存在: {file_path}"
    if path.suffix.lower() != ".pdf":
        return f"错误: 文件 {path.name} 不是 PDF 文件，无法添加 PDF 水印。"

    try:
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import Color
    except ImportError as e:
        return (
            "PDF 加水印工具缺少依赖，请先安装 pypdf 和 reportlab："
            "pip install pypdf reportlab。"
            f"具体错误: {e}"
        )

    try:
        safe_opacity = max(0.0, min(float(opacity), 1.0))
        safe_font_size = max(8, min(int(font_size), 240))
        safe_rotation = int(rotation)
        text = str(watermark_text or "WATERMARK")

        reader = PdfReader(str(path))
        writer = PdfWriter()

        for page in reader.pages:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)

            packet = BytesIO()
            overlay = canvas.Canvas(packet, pagesize=(width, height))
            overlay.saveState()
            if hasattr(overlay, "setFillAlpha"):
                overlay.setFillAlpha(safe_opacity)
            overlay.setFillColor(Color(0.45, 0.45, 0.45, alpha=safe_opacity))
            overlay.setFont("Helvetica-Bold", safe_font_size)
            overlay.translate(width / 2, height / 2)
            overlay.rotate(safe_rotation)
            overlay.drawCentredString(0, 0, text)
            overlay.restoreState()
            overlay.save()

            packet.seek(0)
            watermark_page = PdfReader(packet).pages[0]
            page.merge_page(watermark_page)
            writer.add_page(page)

        if reader.metadata:
            writer.add_metadata(dict(reader.metadata))

        output_name = f"{path.stem}_watermarked_{int(time.time())}.pdf"
        output_path = path.with_name(output_name)
        with open(output_path, "wb") as output:
            writer.write(output)

        return (
            f"已成功添加水印「{text}」。\n"
            f"输出文件: {output_path}\n"
            f"下载链接: /uploads/{output_name}"
        )
    except Exception as e:
        logger.error(f"[Tool:pdf_add_watermark] 执行失败: {e}", exc_info=True)
        return f"PDF 加水印失败: {e}"


@tool
def docx_create(
    content: str,
    title: str = "",
    output_filename: str = "",
) -> str:
    """根据用户提供的内容创建 Word DOCX 文档，并返回生成文件路径。

    Args:
        content: 要写入 Word 文档的正文内容，支持用 Markdown 风格标题和项目符号表达结构
        title: 文档标题，可为空
        output_filename: 输出文件名，必须以 .docx 结尾；可为空自动生成
    """
    logger.info(
        "[Tool:docx_create] title=%s | output=%s | content_len=%s",
        title[:100],
        output_filename,
        len(content or ""),
    )

    if not content or not content.strip():
        return "错误: 生成 Word 文档需要提供正文内容。"

    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError as e:
        return (
            "DOCX 生成工具缺少依赖，请先安装 python-docx："
            "pip install python-docx。"
            f"具体错误: {e}"
        )

    try:
        upload_dir = _get_upload_dir()

        output_name = _safe_output_filename(output_filename, "generated_document", ".docx")
        output_path = upload_dir / output_name
        if output_path.exists():
            output_path = upload_dir / f"{output_path.stem}_{int(time.time())}.docx"
            output_name = output_path.name

        document = Document()
        clean_title = (title or "").strip()
        if clean_title:
            heading = document.add_heading(clean_title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line:
                document.add_paragraph()
                continue

            if line.startswith("### "):
                document.add_heading(line[4:].strip(), level=3)
            elif line.startswith("## "):
                document.add_heading(line[3:].strip(), level=2)
            elif line.startswith("# "):
                document.add_heading(line[2:].strip(), level=1)
            elif line.startswith(("- ", "* ")):
                document.add_paragraph(line[2:].strip(), style="List Bullet")
            elif re.match(r"^\d+[\.)]\s+", line):
                document.add_paragraph(re.sub(r"^\d+[\.)]\s+", "", line).strip(), style="List Number")
            else:
                document.add_paragraph(line)

        document.save(output_path)
        return (
            "已成功生成 Word 文档。\n"
            f"输出文件: {output_path}\n"
            f"下载链接: /uploads/{output_name}"
        )
    except Exception as e:
        logger.error(f"[Tool:docx_create] 执行失败: {e}", exc_info=True)
        return f"Word 文档生成失败: {e}"


@tool
def pptx_create(
    content: str,
    title: str = "",
    output_filename: str = "",
) -> str:
    """根据用户提供的内容创建 PowerPoint PPTX 演示文稿，并返回生成文件路径。

    Args:
        content: 幻灯片内容，建议用 Markdown 标题分隔每页，并用项目符号描述要点
        title: 演示文稿标题，可为空
        output_filename: 输出文件名，必须以 .pptx 结尾；可为空自动生成
    """
    logger.info(
        "[Tool:pptx_create] title=%s | output=%s | content_len=%s",
        title[:100],
        output_filename,
        len(content or ""),
    )

    if not content or not content.strip():
        return "错误: 生成 PPTX 演示文稿需要提供幻灯片内容。"

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as e:
        return (
            "PPTX 生成工具缺少依赖，请先安装 python-pptx："
            "pip install python-pptx。"
            f"具体错误: {e}"
        )

    try:
        upload_dir = _get_upload_dir()
        output_name = _safe_output_filename(output_filename, "generated_presentation", ".pptx")
        output_path = upload_dir / output_name
        if output_path.exists():
            output_path = upload_dir / f"{output_path.stem}_{int(time.time())}.pptx"
            output_name = output_path.name

        slides_data = _parse_slide_sections(content)
        clean_title = (title or "").strip()
        if clean_title and (not slides_data or slides_data[0]["title"] != clean_title):
            slides_data.insert(0, {"title": clean_title, "bullets": []})
        if not slides_data:
            slides_data = [{"title": clean_title or "Presentation", "bullets": [content.strip()]}]

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        prs.core_properties.title = clean_title or slides_data[0]["title"]

        for index, slide_data in enumerate(slides_data[:30]):
            slide = prs.slides.add_slide(blank_layout)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor(248, 250, 252)

            title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(8.7), Inches(0.9))
            title_frame = title_box.text_frame
            title_frame.clear()
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data["title"]
            title_para.font.bold = True
            title_para.font.size = Pt(30 if index else 36)
            title_para.alignment = PP_ALIGN.CENTER if index == 0 else PP_ALIGN.LEFT

            bullets = slide_data.get("bullets") or []
            if bullets:
                body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(8.1), Inches(4.6))
                body_frame = body_box.text_frame
                body_frame.clear()
                body_frame.word_wrap = True
                for bullet_index, bullet in enumerate(bullets[:8]):
                    para = body_frame.paragraphs[0] if bullet_index == 0 else body_frame.add_paragraph()
                    para.text = bullet
                    para.level = 0
                    para.font.size = Pt(20)
                    para.space_after = Pt(8)
            elif index == 0:
                subtitle_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(7.6), Inches(1.0))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.clear()
                subtitle = subtitle_frame.paragraphs[0]
                subtitle.text = "Generated presentation"
                subtitle.font.size = Pt(18)
                subtitle.alignment = PP_ALIGN.CENTER

            footer_box = slide.shapes.add_textbox(Inches(8.75), Inches(6.85), Inches(0.8), Inches(0.25))
            footer_frame = footer_box.text_frame
            footer_frame.clear()
            footer = footer_frame.paragraphs[0]
            footer.text = str(index + 1)
            footer.font.size = Pt(9)
            footer.alignment = PP_ALIGN.RIGHT

        prs.save(output_path)
        return (
            "已成功生成 PowerPoint 演示文稿。\n"
            f"输出文件: {output_path}\n"
            f"下载链接: /uploads/{output_name}"
        )
    except Exception as e:
        logger.error(f"[Tool:pptx_create] 执行失败: {e}", exc_info=True)
        return f"PPTX 演示文稿生成失败: {e}"


@tool
def xlsx_create(
    content: str,
    title: str = "Sheet1",
    output_filename: str = "",
) -> str:
    """根据用户提供的表格内容创建 Excel XLSX 文件，并返回生成文件路径。

    Args:
        content: 表格内容，支持 Markdown 表格、CSV 或 TSV 文本；首行会作为表头
        title: 工作表名称，可为空
        output_filename: 输出文件名，必须以 .xlsx 结尾；可为空自动生成
    """
    logger.info(
        "[Tool:xlsx_create] title=%s | output=%s | content_len=%s",
        title[:100],
        output_filename,
        len(content or ""),
    )

    rows = _parse_table_content(content)
    if not rows:
        return "错误: 生成 Excel 文件需要提供表格内容。"

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError as e:
        return (
            "XLSX 生成工具缺少依赖，请先安装 openpyxl："
            "pip install openpyxl。"
            f"具体错误: {e}"
        )

    try:
        upload_dir = _get_upload_dir()
        output_name = _safe_output_filename(output_filename, "generated_spreadsheet", ".xlsx")
        output_path = upload_dir / output_name
        if output_path.exists():
            output_path = upload_dir / f"{output_path.stem}_{int(time.time())}.xlsx"
            output_name = output_path.name

        wb = Workbook()
        ws = wb.active
        sheet_name = re.sub(r"[\\/*?:\[\]]", "_", (title or "Sheet1").strip())[:31] or "Sheet1"
        ws.title = sheet_name

        max_columns = max(len(row) for row in rows)
        for row in rows:
            row.extend([""] * (max_columns - len(row)))
            ws.append(row)

        header_fill = PatternFill("solid", fgColor="1F4E78")
        header_font = Font(color="FFFFFF", bold=True)
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")

        ws.freeze_panes = "A2"
        ws.auto_filter.ref = ws.dimensions
        for column_cells in ws.columns:
            max_length = max(len(str(cell.value or "")) for cell in column_cells)
            width = min(max(max_length + 2, 10), 45)
            ws.column_dimensions[get_column_letter(column_cells[0].column)].width = width
            for cell in column_cells:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        wb.save(output_path)
        return (
            "已成功生成 Excel 工作簿。\n"
            f"输出文件: {output_path}\n"
            f"下载链接: /uploads/{output_name}"
        )
    except Exception as e:
        logger.error(f"[Tool:xlsx_create] 执行失败: {e}", exc_info=True)
        return f"Excel 文件生成失败: {e}"


# 导出工具列表
GENERAL_CHAT_TOOLS = [
    image_recognition,
    document_summary,
    pdf_add_watermark,
    docx_create,
    pptx_create,
    xlsx_create,
]
